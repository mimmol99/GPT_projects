from docx import Document
from tqdm import tqdm
from pptx import Presentation
import PyPDF2
import os


def file_to_text(file_path):
    _, file_extension = os.path.splitext(file_path)

    text = ""

    file_extensions_available = [".pdf",".doc","docx",".txt",".pptx"]

    if file_extension in file_extensions_available:
        if file_extension == ".pdf":
            text = pdf_to_text(file_path)
        elif file_extension in [".doc", ".docx"]:
            text = docx_to_text(file_path)
        elif file_extension == ".txt":
            text = txt_to_text(file_path)
        elif file_extension == ".pptx":
            text = pptx_to_text(file_path)
    else:
        print(f"File type {file_extension} is not supported. Supported: {file_extensions_available}")
        return None

    return text

def txt_to_text(path):
    with open(path, 'r') as file:
        return file.read()

def pdf_to_text(file_path):
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        num_pages = len(reader.pages)
        text = ''
        
        for page in reader.pages:
            text += page.extract_text()
        
    return text

def docx_to_text(path):
    doc = Document(path)
    fullText = []
    for paragraph in doc.paragraphs:
        fullText.append(paragraph.text)
    return '\n'.join(fullText)

def pptx_to_text(path):
    pres = Presentation(path)
    text = ''
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text