
import os
from pptx import Presentation
from translate import Translator
from tqdm import tqdm
import time

import sys
print(os.path.abspath(os.path.join(os.getcwd(), "..", "..")))
sys.path.append(os.path.abspath(os.path.join(os.getcwd(), "..", "..")))

from SHARED_FILES.directory_request import request_file_path,request_file_paths
from SHARED_FILES.request_gpt import request_gpt

def translate_pptx(file_path, lang="it"):
    translator = Translator(to_lang=lang)
    prs = Presentation(file_path)
    name_prs = os.path.basename(file_path).split(".")[0]
    slides = prs.slides

    with tqdm(total=len(slides), desc="Translating slides") as pbar:
        for slide in slides:
            for shape in slide.shapes:
                # Check if the shape is a textbox
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            # Translate the text and assign it back
                            run.text = request_gpt("Traduci in italiano:" + run.text)
                            time.sleep(1)
            pbar.update(1)

    prs.save(name_prs + "_translated.pptx")


if __name__ == "__main__":
    file_paths = request_file_paths("Select one or more pptx files")

    # Iterate over the file paths
    for file_path in file_paths:
        if file_path.endswith(".pptx"):
            print(file_path, ": translating...")
            translate_pptx(file_path)
        else:
            print(file_path, "is not a pptx, skipping...")


