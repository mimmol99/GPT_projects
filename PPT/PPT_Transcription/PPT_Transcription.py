from docx import Document
from pptx import Presentation
from tqdm import tqdm
import os
import re
from docx.shared import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image as PILImage
from io import BytesIO
import subprocess

import sys

# Get the absolute path to the parent directory
parent_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

# Add it to the Python path
sys.path.insert(0, parent_dir)

from SHARED_FILES.directory_request import request_file_path,request_file_paths
from SHARED_FILES.request_gpt import request_gpt

def convert_wmf_to_png(wmf_bytes):
    process = subprocess.run(
        ['convert', 'wmf:-', 'png:-'],
        input=wmf_bytes,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        check=True,
    )
    return process.stdout

def transcription_pptx(file_path, doc):
    prs = Presentation(file_path)
    name_prs = os.path.basename(file_path).split(".")[0]
    slides = prs.slides

    # Create a new heading for the pptx file
    doc.add_heading(name_prs, level=1)

    # Add tqdm progress bar for the loop
    with tqdm(total=len(slides), desc="Transcribing slides") as pbar:
        for idx, slide in enumerate(slides, 1):
            # Create a new subheading for each slide
            doc.add_heading('Slide ' + str(idx), level=2)
            
            for shape in slide.shapes:
                
                if shape.is_placeholder:
                    # Remove non-XML compatible characters
                    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', shape.text)
                    
                    doc.add_paragraph(text)
            pbar.update(1)  # Update the progress bar

if __name__ == "__main__":
    file_paths = request_file_paths("Select one or more pptx files to be transcripted")
    
    # Create a new Word Document for the transcription
    doc = Document()

    # Iterate over the file paths
    for file_path in file_paths:
        if file_path.endswith(".pptx"):
            print(file_path, ": transcribing...")
            transcription_pptx(file_path, doc)
        else:
            print(file_path, "is not a pptx, skipping...")

    # Save the transcribed text in a .docx file
    doc.save("all_transcribed.docx")

