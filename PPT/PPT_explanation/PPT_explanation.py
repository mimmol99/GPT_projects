from docx import Document
from pptx import Presentation
from tqdm import tqdm
import os
import re
from docx.shared import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image as PILImage
from io import BytesIO
import subprocess
import time
from datetime import datetime

import sys

# Get the absolute path to the parent directory
parent_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

# Add it to the Python path
sys.path.insert(0, parent_dir)

from SHARED_FILES.directory_request import request_file_path, request_file_paths

from SHARED_FILES.request_gpt import request_gpt

def create_directory(directory_name):
    if not os.path.exists(directory_name):
        os.makedirs(directory_name)

def convert_wmf_to_png(wmf_bytes):
    process = subprocess.run(
        ['convert', 'wmf:-', 'png:-'],
        input=wmf_bytes,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        check=True,
    )
    return process.stdout

def transcription_pptx(file_path, doc):
    prs = Presentation(file_path)
    name_prs = os.path.basename(file_path).split(".")[0]
    slides = prs.slides

    # Create a new heading for the pptx file
    doc.add_heading(name_prs, level=1)

    # Add tqdm progress bar for the loop
    with tqdm(total=len(slides), desc="Explaining slides") as pbar:
        for idx, slide in enumerate(slides, 1):
            # Create a new subheading for each slide
            doc.add_heading('Slide ' + str(idx), level=2)

            for shape in slide.shapes:
                if shape.is_placeholder:
                    # Remove non-XML compatible characters
                    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', shape.text)
                    if len(text) > 50:
                        text = request_gpt("spiega e approfondisci questo testo: " + text)
                        time.sleep(5)
                    doc.add_paragraph(text)
            
            pbar.update(1)  # Update the progress bar

if __name__ == "__main__":
    file_paths = request_file_paths("Select one or more pptx files to be explained")

    # Create a new directory based on current time
    if len(file_paths)>0:
        output_directory = datetime.now().strftime("%Y%m%d%H%M%S")
        create_directory(output_directory)

    # Create a new Word Document for the transcription
    

    # Iterate over the file paths
    for file_path in file_paths:
        doc = Document()
        if file_path.endswith(".pptx"):
            print(file_path, ": explaining...")
            transcription_pptx(file_path, doc)

            # Save the transcribed text in a .docx file
            name_prs = os.path.basename(file_path).split(".")[0]
            doc.save(os.path.join(output_directory, name_prs + "_explained.docx"))

            
        else:
            print(file_path, "is not a pptx, skipping...")
        print("Explanation completed!")
		
    # Save the combined transcribed text in a single .docx file
    if len(file_paths) > 1:
        # Combine the document files into a single document
        combined_doc = Document()
        for root, _, files in os.walk(output_directory):
            for file in files:
                if file.endswith(".docx"):
                    file_path = os.path.join(root, file)
                    if file_path != os.path.join(output_directory, "all_explained.docx"):
                        sub_doc = Document(file_path)
                        for element in sub_doc.element.body:
                            combined_doc.element.body.append(element.clone())
        
        # Save the combined document as "all_explained.docx"
        combined_doc.save(os.path.join(output_directory, "all_explained.docx"))

        
